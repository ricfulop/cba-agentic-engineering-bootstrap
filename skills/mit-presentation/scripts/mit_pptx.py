#!/usr/bin/env python3
"""Add slides to the official MIT PowerPoint template.

Uses template layouts. Does not invent a custom black DeckBuilder.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

MIT_RED = RGBColor(163, 31, 52)
DEFAULT_AFFILIATION = (
    "Center for Bits and Atoms, Massachusetts Institute of Technology"
)
LAYOUT_TITLE = 0
LAYOUT_SECTION = 3
LAYOUT_QUOTE = 4


def plugin_root() -> Path:
    return Path(__file__).resolve().parents[3]


def default_template() -> Path:
    return plugin_root() / "templates" / "MIT-PowerPoint-template-Arial-3.pptx"


def _set_run(run, text: str, *, size_pt: float, bold: bool = False, color=None) -> None:
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    latin = rpr.find(qn("a:latin"))
    if latin is None:
        latin = rpr.makeelement(qn("a:latin"), {})
        rpr.insert(0, latin)
    latin.set("typeface", "Arial")


def _fill_placeholder(slide, idx: int, text: str, *, size_pt: float, bold: bool = False) -> None:
    try:
        box = slide.placeholders[idx]
    except KeyError:
        return
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(run, text, size_pt=size_pt, bold=bold)


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    affiliation: str = DEFAULT_AFFILIATION,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _fill_placeholder(slide, 0, title, size_pt=32, bold=True)
    body = subtitle.strip()
    if affiliation:
        body = f"{body}\n{affiliation}" if body else affiliation
    _fill_placeholder(slide, 1, body, size_pt=16)


def add_section_slide(prs: Presentation, heading: str, detail: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    _fill_placeholder(slide, 0, heading, size_pt=28, bold=True)
    if detail:
        _fill_placeholder(slide, 1, detail, size_pt=16)


def add_quote_slide(prs: Presentation, quote: str, attribution: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_QUOTE])
    _fill_placeholder(slide, 0, quote, size_pt=22)
    if attribution:
        _fill_placeholder(slide, 2, attribution, size_pt=14)


def add_content_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    """Reuse a section layout and write Arial body text. No custom chrome."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    _fill_placeholder(slide, 0, heading, size_pt=24, bold=True)
    try:
        box = slide.placeholders[1]
    except KeyError:
        return
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        _set_run(run, item, size_pt=16)
        p.level = 0
    # Silence unused import if a caller wants EMU math later.
    _ = Emu


def build(
    outfile: Path,
    title: str,
    subtitle: str = "",
    affiliation: str = DEFAULT_AFFILIATION,
    sections: list[tuple[str, list[str]]] | None = None,
    template: Path | None = None,
) -> Path:
    src = template or default_template()
    if not src.is_file():
        raise FileNotFoundError(f"MIT template missing: {src}")
    prs = Presentation(str(src))
    add_title_slide(prs, title, subtitle, affiliation)
    for heading, bullets in sections or []:
        add_content_slide(prs, heading, bullets)
    outfile.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(outfile))
    return outfile


def main() -> None:
    parser = argparse.ArgumentParser(description="Official MIT template deck helper")
    parser.add_argument("--title", required=True)
    parser.add_argument("--subtitle", default="")
    parser.add_argument("--affiliation", default=DEFAULT_AFFILIATION)
    parser.add_argument("--outfile", required=True)
    parser.add_argument("--template", default="")
    args = parser.parse_args()
    template = Path(args.template) if args.template else None
    path = build(
        Path(args.outfile),
        args.title,
        subtitle=args.subtitle,
        affiliation=args.affiliation,
        template=template,
    )
    print(path)


if __name__ == "__main__":
    main()
